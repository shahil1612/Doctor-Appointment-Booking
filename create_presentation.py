#!/usr/bin/env python3
"""
Create PowerPoint Presentation for Doctor Appointment Booking System
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    """Create the complete presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Color scheme (light and professional)
    TITLE_COLOR = RGBColor(31, 78, 121)      # Dark Blue
    ACCENT_COLOR = RGBColor(68, 114, 196)    # Medium Blue
    TEXT_COLOR = RGBColor(64, 64, 64)        # Dark Gray
    LIGHT_BG = RGBColor(240, 248, 255)       # Alice Blue
    CODE_BG = RGBColor(248, 248, 248)        # Light Gray
    
    # Slide 1: Title Slide
    slide1 = add_title_slide(prs, TITLE_COLOR, TEXT_COLOR)
    
    # Slide 2: Introduction
    slide2 = add_intro_slide(prs, TITLE_COLOR, TEXT_COLOR, ACCENT_COLOR)
    
    # Slide 3: System Architecture
    slide3 = add_architecture_slide(prs, TITLE_COLOR, TEXT_COLOR, ACCENT_COLOR)
    
    # Slide 4: Technology Stack
    slide4 = add_tech_stack_slide(prs, TITLE_COLOR, TEXT_COLOR, ACCENT_COLOR)
    
    # Slide 5: Database Schema
    slide5 = add_database_slide(prs, TITLE_COLOR, TEXT_COLOR, ACCENT_COLOR)
    
    # Slide 6: Authentication Implementation
    slide6 = add_authentication_slide(prs, TITLE_COLOR, TEXT_COLOR, CODE_BG)
    
    # Slide 7: JWT Token Service
    slide7 = add_jwt_slide(prs, TITLE_COLOR, TEXT_COLOR, CODE_BG)
    
    # Slide 8: Appointment Booking Implementation
    slide8 = add_appointment_booking_slide(prs, TITLE_COLOR, TEXT_COLOR, CODE_BG)
    
    # Slide 9: Payment Integration (Stripe)
    slide9 = add_payment_slide(prs, TITLE_COLOR, TEXT_COLOR, CODE_BG)
    
    # Slide 10: Payment Confirmation Flow
    slide10 = add_payment_confirm_slide(prs, TITLE_COLOR, TEXT_COLOR, CODE_BG)
    
    # Slide 11: Prescription Management
    slide11 = add_prescription_slide(prs, TITLE_COLOR, TEXT_COLOR, CODE_BG)
    
    # Slide 12: Document Upload Implementation
    slide12 = add_document_upload_slide(prs, TITLE_COLOR, TEXT_COLOR, CODE_BG)
    
    # Slide 13: Location-Based Services
    slide13 = add_location_slide(prs, TITLE_COLOR, TEXT_COLOR, CODE_BG)
    
    # Slide 14: Frontend - React Components
    slide14 = add_frontend_slide(prs, TITLE_COLOR, TEXT_COLOR, CODE_BG)
    
    # Slide 15: Redux State Management
    slide15 = add_redux_slide(prs, TITLE_COLOR, TEXT_COLOR, CODE_BG)
    
    # Slide 16: API Service Layer
    slide16 = add_api_service_slide(prs, TITLE_COLOR, TEXT_COLOR, CODE_BG)
    
    # Slide 17: Security Implementation
    slide17 = add_security_slide(prs, TITLE_COLOR, TEXT_COLOR, ACCENT_COLOR)
    
    # Slide 18: Middleware Pipeline
    slide18 = add_middleware_slide(prs, TITLE_COLOR, TEXT_COLOR, CODE_BG)
    
    # Slide 19: Repository Pattern
    slide19 = add_repository_slide(prs, TITLE_COLOR, TEXT_COLOR, CODE_BG)
    
    # Slide 20: API Endpoints Overview
    slide20 = add_api_endpoints_slide(prs, TITLE_COLOR, TEXT_COLOR, ACCENT_COLOR)
    
    # Slide 21: Key Features Summary
    slide21 = add_features_summary_slide(prs, TITLE_COLOR, TEXT_COLOR, ACCENT_COLOR)
    
    # Slide 22: Screenshots (Placeholder)
    slide22 = add_screenshots_slide(prs, TITLE_COLOR, TEXT_COLOR)
    
    # Slide 23: Conclusion
    slide23 = add_conclusion_slide(prs, TITLE_COLOR, TEXT_COLOR, ACCENT_COLOR)
    
    # Slide 24: References
    slide24 = add_references_slide(prs, TITLE_COLOR, TEXT_COLOR)
    
    # Save presentation
    prs.save('Sehat_Project_Presentation.pptx')
    print("✅ Presentation created successfully: Sehat_Project_Presentation.pptx")
    print("📊 Total slides: 24")
    print("📝 Next steps: Fill in student/guide details on Slide 1 and convert to PDF")

def add_title_slide(prs, title_color, text_color):
    """Slide 1: Title Slide with placeholders"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add decorative shape at top
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0), Inches(0), Inches(10), Inches(1.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(68, 114, 196)
    shape.line.color.rgb = RGBColor(68, 114, 196)
    
    # Project Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "SEHAT - Doctor Appointment Booking System"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(1.7), Inches(8), Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "A Full-Stack Healthcare Appointment Management System"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(18)
    subtitle_para.font.italic = True
    subtitle_para.font.color.rgb = text_color
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Student Names Section
    y_pos = 2.5
    add_info_section(slide, "Student Names:", "[Student Name 1]\n[Student Name 2]\n[Student Name 3]", 
                     Inches(1), Inches(y_pos), text_color)
    
    # Company Name
    y_pos = 3.5
    add_info_section(slide, "Company Name:", "[Company/Organization Name]", 
                     Inches(1), Inches(y_pos), text_color)
    
    # Guides Section
    y_pos = 4.3
    guides_box = slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(8), Inches(2))
    guides_frame = guides_box.text_frame
    guides_frame.word_wrap = True
    
    p1 = guides_frame.paragraphs[0]
    p1.text = "External Guide:"
    p1.font.size = Pt(14)
    p1.font.bold = True
    p1.font.color.rgb = title_color
    
    p2 = guides_frame.add_paragraph()
    p2.text = "[Guide Name], [Designation]\n[Company Name]"
    p2.font.size = Pt(12)
    p2.font.color.rgb = text_color
    p2.space_before = Pt(3)
    
    p3 = guides_frame.add_paragraph()
    p3.text = "\nInternal Guide:"
    p3.font.size = Pt(14)
    p3.font.bold = True
    p3.font.color.rgb = title_color
    p3.space_before = Pt(8)
    
    p4 = guides_frame.add_paragraph()
    p4.text = "[Guide Name], [Designation]\n[Institution Name]"
    p4.font.size = Pt(12)
    p4.font.color.rgb = text_color
    p4.space_before = Pt(3)
    
    # Footer with date
    footer_box = slide.shapes.add_textbox(Inches(1), Inches(6.8), Inches(8), Inches(0.5))
    footer_frame = footer_box.text_frame
    footer_frame.text = "April 2026"
    footer_para = footer_frame.paragraphs[0]
    footer_para.font.size = Pt(10)
    footer_para.font.color.rgb = text_color
    footer_para.alignment = PP_ALIGN.CENTER
    
    return slide

def add_info_section(slide, label, value, left, top, text_color):
    """Helper to add information sections"""
    box = slide.shapes.add_textbox(left, top, Inches(8), Inches(0.8))
    frame = box.text_frame
    frame.word_wrap = True
    
    p1 = frame.paragraphs[0]
    p1.text = label
    p1.font.size = Pt(14)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(31, 78, 121)
    
    p2 = frame.add_paragraph()
    p2.text = value
    p2.font.size = Pt(12)
    p2.font.color.rgb = text_color
    p2.space_before = Pt(3)

def add_intro_slide(prs, title_color, text_color, accent_color):
    """Slide 2: Introduction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Introduction", title_color)
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
    frame = content_box.text_frame
    frame.word_wrap = True
    
    sections = [
        ("Project Overview:", "Sehat is a comprehensive healthcare appointment management system that connects patients with doctors through a modern, secure platform."),
        ("Problem Statement:", "• Inefficient appointment booking processes\n• Lack of integrated payment systems\n• Poor digital health record management\n• Limited access to doctor information"),
        ("Solution:", "A full-stack web application providing:\n• Easy doctor discovery with location-based search\n• Seamless appointment booking with real-time availability\n• Secure payment processing via Stripe\n• Digital prescription and document management\n• Multi-clinic support for healthcare providers"),
        ("Project Scope:", "End-to-end implementation from database design to user interface, including backend APIs, frontend dashboard, payment integration, and security features.")
    ]
    
    y_offset = 0
    for label, text in sections:
        p = frame.paragraphs[0] if y_offset == 0 else frame.add_paragraph()
        p.text = label
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = accent_color
        if y_offset > 0:
            p.space_before = Pt(12)
        
        p2 = frame.add_paragraph()
        p2.text = text
        p2.font.size = Pt(11)
        p2.font.color.rgb = text_color
        p2.space_before = Pt(3)
        y_offset += 1
    
    return slide

def add_architecture_slide(prs, title_color, text_color, accent_color):
    """Slide 3: System Architecture"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "System Architecture", title_color)
    
    # Architecture diagram (text-based)
    arch_box = slide.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(7), Inches(4.5))
    frame = arch_box.text_frame
    frame.word_wrap = True
    
    p1 = frame.paragraphs[0]
    p1.text = "Layered Architecture (N-Tier Pattern)"
    p1.font.size = Pt(14)
    p1.font.bold = True
    p1.font.color.rgb = accent_color
    p1.alignment = PP_ALIGN.CENTER
    
    architecture_text = """
┌─────────────────────────────────┐
│   Presentation Layer (React)    │
│   • Components & Pages          │
│   • Redux State Management      │
└────────────┬────────────────────┘
             │ HTTP/REST
┌────────────▼────────────────────┐
│   API Layer (Controllers)       │
│   • RESTful Endpoints           │
│   • Request/Response Handling   │
└────────────┬────────────────────┘
             │
┌────────────▼────────────────────┐
│   Business Logic (Services)     │
│   • Validation & Processing     │
│   • Business Rules              │
└────────────┬────────────────────┘
             │
┌────────────▼────────────────────┐
│   Data Access (Repositories)    │
│   • Entity Framework Core       │
│   • CRUD Operations             │
└────────────┬────────────────────┘
             │
┌────────────▼────────────────────┐
│   Database Layer (MySQL)        │
│   • 11 Tables                   │
│   • Relationships & Constraints │
└─────────────────────────────────┘
"""
    
    p2 = frame.add_paragraph()
    p2.text = architecture_text
    p2.font.name = 'Courier New'
    p2.font.size = Pt(9)
    p2.font.color.rgb = text_color
    p2.space_before = Pt(10)
    
    # Design Patterns
    patterns_box = slide.shapes.add_textbox(Inches(1), Inches(6.3), Inches(8), Inches(0.8))
    patterns_frame = patterns_box.text_frame
    p = patterns_frame.paragraphs[0]
    p.text = "Design Patterns: Repository Pattern • Service Pattern • Dependency Injection • DTO Pattern • Middleware Pipeline"
    p.font.size = Pt(10)
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_tech_stack_slide(prs, title_color, text_color, accent_color):
    """Slide 4: Technology Stack"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Technology Stack", title_color)
    
    # Backend
    backend_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(4), Inches(4.5))
    frame = backend_box.text_frame
    frame.word_wrap = True
    
    p1 = frame.paragraphs[0]
    p1.text = "Backend Technologies"
    p1.font.size = Pt(16)
    p1.font.bold = True
    p1.font.color.rgb = accent_color
    
    backend_tech = [
        "• ASP.NET Core 8.0 Web API",
        "• C# 12 (Nullable enabled)",
        "• Entity Framework Core 8.0",
        "• MySQL 8.0 (Pomelo Provider)",
        "• JWT Bearer Authentication",
        "• BCrypt.Net (Password Hashing)",
        "• Stripe.NET SDK 51.0",
        "• NLog (Structured Logging)",
        "• Swagger/OpenAPI",
        "• Dependency Injection"
    ]
    
    for tech in backend_tech:
        p = frame.add_paragraph()
        p.text = tech
        p.font.size = Pt(11)
        p.font.color.rgb = text_color
        p.space_before = Pt(3)
        p.level = 0
    
    # Frontend
    frontend_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.8), Inches(4), Inches(4.5))
    frame2 = frontend_box.text_frame
    frame2.word_wrap = True
    
    p1 = frame2.paragraphs[0]
    p1.text = "Frontend Technologies"
    p1.font.size = Pt(16)
    p1.font.bold = True
    p1.font.color.rgb = accent_color
    
    frontend_tech = [
        "• React 19.2",
        "• Vite 7.3 (Build Tool)",
        "• Tailwind CSS 4.2",
        "• Redux Toolkit 2.11",
        "• React Router DOM 7",
        "• Framer Motion (Animations)",
        "• React Leaflet (Maps)",
        "• Stripe React Components",
        "• React Hot Toast",
        "• Lucide React (Icons)"
    ]
    
    for tech in frontend_tech:
        p = frame2.add_paragraph()
        p.text = tech
        p.font.size = Pt(11)
        p.font.color.rgb = text_color
        p.space_before = Pt(3)
        p.level = 0
    
    return slide

def add_database_slide(prs, title_color, text_color, accent_color):
    """Slide 5: Database Schema"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Database Schema", title_color)
    
    # Table listing
    tables_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(4.8))
    frame = tables_box.text_frame
    frame.word_wrap = True
    
    p1 = frame.paragraphs[0]
    p1.text = "Database: MySQL 8.0 | Total Tables: 11 | ORM: Entity Framework Core"
    p1.font.size = Pt(12)
    p1.font.bold = True
    p1.font.color.rgb = accent_color
    p1.alignment = PP_ALIGN.CENTER
    
    tables = [
        ("TBL01 - users", "Master user accounts (email, password_hash, user_type)"),
        ("TBL02 - patients", "Patient health profiles (blood_group, allergies, emergency_contact)"),
        ("TBL03 - doctors", "Doctor credentials (specialization, license, consultation_fee)"),
        ("TBL04 - appointments", "Appointment records (status, reason, appointment_time)"),
        ("TBL05 - clinics", "Healthcare facilities (address, coordinates, phone)"),
        ("TBL06 - doctor_clinics", "Doctor-Clinic many-to-many mapping"),
        ("TBL07 - appointment_slots", "Available time slots (slot_start, is_booked)"),
        ("TBL08 - prescriptions", "Digital prescriptions (diagnosis, medications, dosage)"),
        ("TBL09 - medical_documents", "Patient documents (file_path, file_size, mime_type)"),
        ("TBL10 - payments", "Stripe transactions (amount_cents, status, payment_intent_id)"),
        ("TBL11 - refunds", "Payment refunds (refund_amount, refund_status)")
    ]
    
    for table_name, desc in tables:
        p = frame.add_paragraph()
        p.text = f"{table_name}: {desc}"
        p.font.size = Pt(10)
        p.font.color.rgb = text_color
        p.space_before = Pt(6)
        run = p.runs[0]
        # Bold the table name part
        parts = p.text.split(': ', 1)
        p.text = parts[0]
        p.font.bold = True
        p.font.color.rgb = accent_color
        run2 = p.add_run()
        run2.text = f": {parts[1]}"
        run2.font.bold = False
        run2.font.color.rgb = text_color
    
    # Relationships note
    rel_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.6))
    rel_frame = rel_box.text_frame
    p = rel_frame.paragraphs[0]
    p.text = "Key Relationships: User ↔ Patient/Doctor (1:1) • Appointments → Payments (1:1) • Doctors ↔ Clinics (M:N)"
    p.font.size = Pt(10)
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_authentication_slide(prs, title_color, text_color, code_bg):
    """Slide 6: Authentication Implementation"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Implementation: Authentication Service", title_color)
    
    # Description
    desc_box = slide.shapes.add_textbox(Inches(1), Inches(1.6), Inches(8), Inches(0.6))
    desc_frame = desc_box.text_frame
    p = desc_frame.paragraphs[0]
    p.text = "Secure user registration and login with BCrypt password hashing and JWT token generation"
    p.font.size = Pt(11)
    p.font.color.rgb = text_color
    
    # Code snippet
    code = """// AuthService.cs - Password Hashing & User Creation

public async Task<LoginResponse> SignupPresaveAsync(SignupRequest request)
{
    // Check if email already exists
    var existingUser = await _userRepository
        .FindByEmailAsync(request.Email);
    if (existingUser != null)
        throw new AppException("Email already registered", 400);

    // Hash password using BCrypt
    string passwordHash = BCrypt.Net.BCrypt.HashPassword(
        request.Password, 
        workFactor: 4
    );

    // Create user entity
    var user = new TBL01
    {
        L01F02 = (int)request.UserType,  // PATIENT or DOCTOR
        L01F03 = request.FullName,
        L01F04 = request.Email,
        L01F05 = passwordHash,           // Hashed password
        L01F06 = request.PhoneNumber,
        L01F07 = DateTime.UtcNow
    };

    return user;
}"""
    
    add_code_box(slide, code, Inches(1), Inches(2.3), Inches(8), Inches(4), code_bg, text_color)
    
    # Key points
    points_box = slide.shapes.add_textbox(Inches(1), Inches(6.4), Inches(8), Inches(0.6))
    points_frame = points_box.text_frame
    p = points_frame.paragraphs[0]
    p.text = "✓ BCrypt hashing with work factor 4  •  ✓ Email uniqueness validation  •  ✓ UTC timestamp tracking"
    p.font.size = Pt(10)
    p.font.color.rgb = text_color
    
    return slide

def add_jwt_slide(prs, title_color, text_color, code_bg):
    """Slide 7: JWT Token Service"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Implementation: JWT Token Generation", title_color)
    
    # Description
    desc_box = slide.shapes.add_textbox(Inches(1), Inches(1.6), Inches(8), Inches(0.5))
    desc_frame = desc_box.text_frame
    p = desc_frame.paragraphs[0]
    p.text = "Generate JWT tokens with user claims for authentication and authorization"
    p.font.size = Pt(11)
    p.font.color.rgb = text_color
    
    # Code snippet
    code = """// TokenService.cs - JWT Token Generation

public string GenerateToken(TBL01 user)
{
    var key = new SymmetricSecurityKey(
        Convert.FromBase64String(_jwtSecret)
    );
    var credentials = new SigningCredentials(
        key, 
        SecurityAlgorithms.HmacSha256
    );

    // Define claims
    var claims = new[]
    {
        new Claim(ClaimTypes.NameIdentifier, user.L01F01.ToString()),
        new Claim(ClaimTypes.Role, ((UserType)user.L01F02).ToString())
    };

    // Create JWT token
    var token = new JwtSecurityToken(
        claims: claims,
        expires: DateTime.UtcNow.AddHours(2),  // 2-hour expiration
        signingCredentials: credentials
    );

    return new JwtSecurityTokenHandler().WriteToken(token);
}"""
    
    add_code_box(slide, code, Inches(1), Inches(2.2), Inches(8), Inches(3.8), code_bg, text_color)
    
    # JWT structure
    jwt_box = slide.shapes.add_textbox(Inches(1), Inches(6.1), Inches(8), Inches(0.8))
    jwt_frame = jwt_box.text_frame
    p = jwt_frame.paragraphs[0]
    p.text = "JWT Structure: Header (Algorithm) • Payload (User ID, Role) • Signature (HMAC-SHA256)"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = RGBColor(31, 78, 121)
    
    p2 = jwt_frame.add_paragraph()
    p2.text = "Token expires after 2 hours • Stored in localStorage on client • Sent as Bearer token in Authorization header"
    p2.font.size = Pt(9)
    p2.font.color.rgb = text_color
    
    return slide

def add_appointment_booking_slide(prs, title_color, text_color, code_bg):
    """Slide 8: Appointment Booking"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Implementation: Appointment Booking", title_color)
    
    desc_box = slide.shapes.add_textbox(Inches(1), Inches(1.6), Inches(8), Inches(0.5))
    desc_frame = desc_box.text_frame
    p = desc_frame.paragraphs[0]
    p.text = "Multi-step workflow for creating appointments with validation and slot management"
    p.font.size = Pt(11)
    p.font.color.rgb = text_color
    
    code = """// AppointmentService.cs - Create Appointment Workflow

public async Task CreateAppointmentValidateAsync()
{
    var slot = await _appointmentSlotRepository
        .FindByIdAsync(_createState.Request.SlotId);
    
    if (slot == null)
        throw new AppException("Slot not found", 404);
    
    if (slot.L07F06)  // is_booked flag
        throw new AppException("Slot already booked", 400);
    
    // Verify clinic association
    var doctorClinic = await _doctorClinicRepository
        .FindByDoctorAndClinicAsync(
            slot.L07F02,  // doctor_user_id
            _createState.Request.ClinicId
        );
    
    if (doctorClinic == null)
        throw new AppException("Invalid clinic for doctor", 400);
}

public async Task<AppointmentResponse> CreateAppointmentSaveAsync()
{
    // Create appointment with PENDING status
    var appointment = new TBL04
    {
        L04F02 = _createState.PatientUserId,
        L04F03 = _createState.Slot.L07F02,  // doctor_user_id
        L04F04 = _createState.Request.SlotId,
        L04F05 = _createState.Slot.L07F04,  // slot_start_utc
        L04F06 = (int)AppointmentStatus.PENDING,
        L04F07 = _createState.Request.Reason,
        L04F10 = _createState.Request.ClinicId
    };

    await _appointmentRepository.AddAsync(appointment);
    
    // Mark slot as booked
    _createState.Slot.L07F06 = true;
    await _appointmentSlotRepository.UpdateAsync(_createState.Slot);
    
    return _mapper.Map<AppointmentResponse>(appointment);
}"""
    
    add_code_box(slide, code, Inches(0.7), Inches(2.2), Inches(8.6), Inches(4.3), code_bg, text_color)
    
    return slide

def add_payment_slide(prs, title_color, text_color, code_bg):
    """Slide 9: Payment Integration"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Implementation: Stripe Payment Integration", title_color)
    
    desc_box = slide.shapes.add_textbox(Inches(1), Inches(1.6), Inches(8), Inches(0.5))
    desc_frame = desc_box.text_frame
    p = desc_frame.paragraphs[0]
    p.text = "Create payment intent and handle Stripe payment processing"
    p.font.size = Pt(11)
    p.font.color.rgb = text_color
    
    code = """// PaymentService.cs - Create Stripe Payment Intent

public async Task<PaymentIntentResponse> CreatePaymentIntentAsync(
    int patientUserId, 
    CreatePaymentIntentRequest request)
{
    // Get slot and doctor-clinic details
    var slot = await _appointmentRepository.FindSlotByIdAsync(request.SlotId);
    var doctorClinic = await _doctorClinicRepository
        .FindByDoctorAndClinicAsync(slot.L07F02, slot.L07F03);

    // Get consultation fee (clinic-specific or doctor default)
    long amountCents = doctorClinic?.L06F04 ?? 
                       (await _doctorRepository.FindByUserIdAsync(slot.L07F02))
                       .L03F09;

    // Create Stripe payment intent
    var options = new PaymentIntentCreateOptions
    {
        Amount = amountCents,
        Currency = _stripeSettings.Currency,  // "inr"
        Metadata = new Dictionary<string, string>
        {
            { "patient_user_id", patientUserId.ToString() },
            { "doctor_user_id", slot.L07F02.ToString() },
            { "slot_id", request.SlotId.ToString() },
            { "reason", request.Reason }
        }
    };

    var service = new PaymentIntentService();
    var paymentIntent = await service.CreateAsync(options);

    // Store payment record with PENDING status
    var payment = new TBL10
    {
        L10F03 = paymentIntent.Id,        // stripe_payment_intent_id
        L10F04 = amountCents,
        L10F05 = _stripeSettings.Currency,
        L10F06 = (int)PaymentStatus.PENDING
    };
    
    await _paymentRepository.AddAsync(payment);

    return new PaymentIntentResponse
    {
        ClientSecret = paymentIntent.ClientSecret,
        AmountCents = amountCents
    };
}"""
    
    add_code_box(slide, code, Inches(0.5), Inches(2.2), Inches(9), Inches(4.5), code_bg, text_color)
    
    return slide

def add_payment_confirm_slide(prs, title_color, text_color, code_bg):
    """Slide 10: Payment Confirmation"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Implementation: Payment Confirmation Flow", title_color)
    
    desc_box = slide.shapes.add_textbox(Inches(1), Inches(1.6), Inches(8), Inches(0.5))
    desc_frame = desc_box.text_frame
    p = desc_frame.paragraphs[0]
    p.text = "Verify Stripe payment and create appointment after successful payment"
    p.font.size = Pt(11)
    p.font.color.rgb = text_color
    
    code = """// PaymentService.cs - Confirm Payment & Create Appointment

public async Task<PaymentConfirmResponse> ConfirmPaymentAsync(
    int patientUserId, 
    ConfirmPaymentRequest request)
{
    // Retrieve payment intent from Stripe
    var service = new PaymentIntentService();
    var paymentIntent = await service.GetAsync(request.PaymentIntentId);

    if (paymentIntent.Status != "succeeded")
        throw new AppException("Payment not successful", 400);

    // Update payment record
    var payment = await _paymentRepository
        .FindByStripeIdAsync(request.PaymentIntentId);
    payment.L10F06 = (int)PaymentStatus.SUCCEEDED;
    payment.L10F07 = paymentIntent.PaymentMethodId;
    await _paymentRepository.UpdateAsync(payment);

    // Create appointment using appointment service
    var appointmentRequest = new CreateAppointmentRequest
    {
        SlotId = request.SlotId,
        Reason = request.Reason,
        ClinicId = /* from slot */
    };

    await _appointmentService.CreateAppointmentPresaveAsync(
        patientUserId, 
        appointmentRequest
    );
    await _appointmentService.CreateAppointmentValidateAsync();
    var appointment = await _appointmentService.CreateAppointmentSaveAsync();

    // Link payment to appointment
    payment.L10F02 = appointment.AppointmentId;
    await _paymentRepository.UpdateAsync(payment);

    return new PaymentConfirmResponse
    {
        AppointmentId = appointment.AppointmentId,
        PaymentId = payment.L10F01,
        Status = "SUCCEEDED"
    };
}"""
    
    add_code_box(slide, code, Inches(0.5), Inches(2.2), Inches(9), Inches(4.5), code_bg, text_color)
    
    return slide

def add_prescription_slide(prs, title_color, text_color, code_bg):
    """Slide 11: Prescription Management"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Implementation: Digital Prescription System", title_color)
    
    desc_box = slide.shapes.add_textbox(Inches(1), Inches(1.6), Inches(8), Inches(0.5))
    desc_frame = desc_box.text_frame
    p = desc_frame.paragraphs[0]
    p.text = "Create and manage digital prescriptions post-appointment"
    p.font.size = Pt(11)
    p.font.color.rgb = text_color
    
    code = """// PrescriptionService.cs - Create Prescription

public async Task<PrescriptionResponse> CreatePrescriptionAsync(
    int doctorUserId, 
    CreatePrescriptionRequest request)
{
    // Verify appointment exists and is completed
    var appointment = await _appointmentRepository
        .FindByIdAsync(request.AppointmentId);
    
    if (appointment == null)
        throw new AppException("Appointment not found", 404);
    
    if (appointment.L04F03 != doctorUserId)
        throw new AppException("Not authorized", 403);
    
    if (appointment.L04F06 != (int)AppointmentStatus.COMPLETED)
        throw new AppException("Appointment not completed", 400);

    // Create prescription entity
    var prescription = new TBL08
    {
        L08F02 = request.AppointmentId,
        L08F03 = request.Diagnosis,
        L08F04 = request.Medications,
        L08F05 = request.Dosage,
        L08F06 = request.Frequency,
        L08F07 = request.Duration,
        L08F08 = request.Instructions,
        L08F09 = DateTime.UtcNow
    };

    await _prescriptionRepository.AddAsync(prescription);

    return _mapper.Map<PrescriptionResponse>(prescription);
}"""
    
    add_code_box(slide, code, Inches(0.8), Inches(2.2), Inches(8.4), Inches(4.3), code_bg, text_color)
    
    # Key points
    points_box = slide.shapes.add_textbox(Inches(1), Inches(6.6), Inches(8), Inches(0.4))
    points_frame = points_box.text_frame
    p = points_frame.paragraphs[0]
    p.text = "✓ Authorization check (doctor owns appointment)  •  ✓ Status validation (must be COMPLETED)  •  ✓ Structured fields"
    p.font.size = Pt(9)
    p.font.color.rgb = text_color
    
    return slide

def add_document_upload_slide(prs, title_color, text_color, code_bg):
    """Slide 12: Document Upload"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Implementation: Medical Document Upload", title_color)
    
    desc_box = slide.shapes.add_textbox(Inches(1), Inches(1.6), Inches(8), Inches(0.5))
    desc_frame = desc_box.text_frame
    p = desc_frame.paragraphs[0]
    p.text = "Secure file upload with validation and storage management"
    p.font.size = Pt(11)
    p.font.color.rgb = text_color
    
    code = """// MedicalDocumentService.cs - Upload Document

public async Task<MedicalDocumentResponse> UploadDocumentAsync(
    int patientUserId, 
    UploadMedicalDocumentRequest request)
{
    // Validate file size (10 MB limit)
    const long maxSize = 10 * 1024 * 1024;
    if (request.File.Length > maxSize)
        throw new AppException("File too large (max 10MB)", 400);

    // Validate file type
    var allowedTypes = new[] { "application/pdf", "image/jpeg", 
                               "image/png", "application/msword" };
    if (!allowedTypes.Contains(request.File.ContentType))
        throw new AppException("Invalid file type", 400);

    // Generate unique filename
    var fileName = $"{Guid.NewGuid()}_{request.File.FileName}";
    var uploadPath = Path.Combine(_webHostEnvironment.WebRootPath, 
                                   "uploads", fileName);

    // Save file to disk
    using (var stream = new FileStream(uploadPath, FileMode.Create))
    {
        await request.File.CopyToAsync(stream);
    }

    // Store metadata in database
    var document = new TBL09
    {
        L09F02 = patientUserId,
        L09F03 = request.Title,
        L09F04 = request.Description,
        L09F05 = fileName,
        L09F06 = uploadPath,
        L09F07 = (int)request.File.Length,
        L09F08 = request.File.ContentType,
        L09F09 = DateTime.UtcNow
    };

    await _medicalDocumentRepository.AddAsync(document);

    return _mapper.Map<MedicalDocumentResponse>(document);
}"""
    
    add_code_box(slide, code, Inches(0.5), Inches(2.2), Inches(9), Inches(4.5), code_bg, text_color)
    
    return slide

def add_location_slide(prs, title_color, text_color, code_bg):
    """Slide 13: Location-Based Services"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Implementation: Location-Based Doctor Search", title_color)
    
    desc_box = slide.shapes.add_textbox(Inches(1), Inches(1.6), Inches(8), Inches(0.5))
    desc_frame = desc_box.text_frame
    p = desc_frame.paragraphs[0]
    p.text = "Calculate distance between patient and doctors using Haversine formula"
    p.font.size = Pt(11)
    p.font.color.rgb = text_color
    
    code = """// AppointmentService.cs - Distance Calculation (Haversine Formula)

public async Task<List<AvailableDoctorResponse>> GetAvailableDoctorsAsync(
    int patientUserId, 
    double? patientLat, 
    double? patientLong)
{
    var doctors = await _doctorRepository.GetAllWithClinicsAsync();
    var availableDoctors = new List<AvailableDoctorResponse>();

    foreach (var doctor in doctors)
    {
        foreach (var clinic in doctor.DoctorClinics)
        {
            double distance = 0;
            
            if (patientLat.HasValue && patientLong.HasValue)
            {
                // Haversine formula for distance calculation
                const double EarthRadiusKm = 6371.0;
                
                double dLat = ToRadians(clinic.Clinic.L05F08 - patientLat.Value);
                double dLon = ToRadians(clinic.Clinic.L05F09 - patientLong.Value);
                
                double a = Math.Sin(dLat / 2) * Math.Sin(dLat / 2) +
                          Math.Cos(ToRadians(patientLat.Value)) * 
                          Math.Cos(ToRadians(clinic.Clinic.L05F08)) *
                          Math.Sin(dLon / 2) * Math.Sin(dLon / 2);
                
                double c = 2 * Math.Atan2(Math.Sqrt(a), Math.Sqrt(1 - a));
                distance = EarthRadiusKm * c;
            }

            availableDoctors.Add(new AvailableDoctorResponse
            {
                DoctorUserId = doctor.L03F02,
                FullName = doctor.User.L01F03,
                Specialization = doctor.L03F03,
                Distance = Math.Round(distance, 2),  // in kilometers
                ClinicName = clinic.Clinic.L05F02
            });
        }
    }

    return availableDoctors.OrderBy(d => d.Distance).ToList();
}

private double ToRadians(double degrees) => degrees * Math.PI / 180.0;"""
    
    add_code_box(slide, code, Inches(0.4), Inches(2.2), Inches(9.2), Inches(4.5), code_bg, text_color)
    
    return slide

def add_frontend_slide(prs, title_color, text_color, code_bg):
    """Slide 14: Frontend React Components"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Implementation: React Frontend Components", title_color)
    
    desc_box = slide.shapes.add_textbox(Inches(1), Inches(1.6), Inches(8), Inches(0.5))
    desc_frame = desc_box.text_frame
    p = desc_frame.paragraphs[0]
    p.text = "Component-based architecture with React hooks and routing"
    p.font.size = Pt(11)
    p.font.color.rgb = text_color
    
    code = """// App.jsx - Main Application Component with Routing

import { BrowserRouter, Routes, Route } from 'react-router-dom';
import { Provider } from 'react-redux';
import { store } from './store';
import LandingPage from './components/landing/LandingPage';
import Login from './components/auth/Login';
import Signup from './components/auth/Signup';
import PatientDashboard from './components/dashboard/patient/PatientDashboard';
import DoctorDashboard from './components/dashboard/doctor/DoctorDashboard';
import ProtectedRoute from './components/auth/ProtectedRoute';

function App() {
  return (
    <Provider store={store}>
      <BrowserRouter>
        <Routes>
          <Route path="/" element={<LandingPage />} />
          <Route path="/login" element={<Login />} />
          <Route path="/signup" element={<Signup />} />
          
          {/* Protected Routes */}
          <Route 
            path="/patient/*" 
            element={
              <ProtectedRoute allowedRoles={['PATIENT']}>
                <PatientDashboard />
              </ProtectedRoute>
            } 
          />
          
          <Route 
            path="/doctor/*" 
            element={
              <ProtectedRoute allowedRoles={['DOCTOR']}>
                <DoctorDashboard />
              </ProtectedRoute>
            } 
          />
        </Routes>
      </BrowserRouter>
    </Provider>
  );
}

export default App;"""
    
    add_code_box(slide, code, Inches(0.7), Inches(2.2), Inches(8.6), Inches(4.3), code_bg, text_color)
    
    return slide

def add_redux_slide(prs, title_color, text_color, code_bg):
    """Slide 15: Redux State Management"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Implementation: Redux State Management", title_color)
    
    desc_box = slide.shapes.add_textbox(Inches(1), Inches(1.6), Inches(8), Inches(0.5))
    desc_frame = desc_box.text_frame
    p = desc_frame.paragraphs[0]
    p.text = "Centralized state management for authentication and user data"
    p.font.size = Pt(11)
    p.font.color.rgb = text_color
    
    code = """// store/authSlice.js - Redux Authentication Slice

import { createSlice } from '@reduxjs/toolkit';

const initialState = {
  isAuthenticated: false,
  token: localStorage.getItem('token') || null,
  user: JSON.parse(localStorage.getItem('user')) || null,
  userType: localStorage.getItem('userType') || null,
};

const authSlice = createSlice({
  name: 'auth',
  initialState,
  reducers: {
    loginSuccess: (state, action) => {
      state.isAuthenticated = true;
      state.token = action.payload.token;
      state.user = action.payload.user;
      state.userType = action.payload.userType;
      
      // Persist to localStorage
      localStorage.setItem('token', action.payload.token);
      localStorage.setItem('user', JSON.stringify(action.payload.user));
      localStorage.setItem('userType', action.payload.userType);
    },
    
    logout: (state) => {
      state.isAuthenticated = false;
      state.token = null;
      state.user = null;
      state.userType = null;
      
      // Clear localStorage
      localStorage.removeItem('token');
      localStorage.removeItem('user');
      localStorage.removeItem('userType');
    },
    
    updateUser: (state, action) => {
      state.user = { ...state.user, ...action.payload };
      localStorage.setItem('user', JSON.stringify(state.user));
    }
  }
});

export const { loginSuccess, logout, updateUser } = authSlice.actions;
export default authSlice.reducer;"""
    
    add_code_box(slide, code, Inches(0.6), Inches(2.2), Inches(8.8), Inches(4.4), code_bg, text_color)
    
    return slide

def add_api_service_slide(prs, title_color, text_color, code_bg):
    """Slide 16: API Service Layer"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Implementation: API Service Layer", title_color)
    
    desc_box = slide.shapes.add_textbox(Inches(1), Inches(1.6), Inches(8), Inches(0.5))
    desc_frame = desc_box.text_frame
    p = desc_frame.paragraphs[0]
    p.text = "Centralized API client with automatic token injection"
    p.font.size = Pt(11)
    p.font.color.rgb = text_color
    
    code = """// services/api.js - API Client Service

const API_BASE_URL = import.meta.env.VITE_API_BASE_URL;

// Helper to get auth token
const getAuthToken = () => localStorage.getItem('token');

// Generic API request function
const apiRequest = async (endpoint, method = 'GET', body = null) => {
  const headers = {
    'Content-Type': 'application/json',
  };

  // Add authorization header if token exists
  const token = getAuthToken();
  if (token) {
    headers['Authorization'] = `Bearer ${token}`;
  }

  const config = {
    method,
    headers,
  };

  if (body && method !== 'GET') {
    config.body = JSON.stringify(body);
  }

  const response = await fetch(`${API_BASE_URL}${endpoint}`, config);
  
  if (!response.ok) {
    const error = await response.json();
    throw new Error(error.message || 'API request failed');
  }

  return response.json();
};

// Authentication APIs
export const authAPI = {
  login: (credentials) => apiRequest('/auth/login', 'POST', credentials),
  signup: (userData) => apiRequest('/auth/signup', 'POST', userData),
};

// Appointment APIs
export const appointmentAPI = {
  getAvailableDoctors: () => apiRequest('/appointments/doctors/available'),
  getAvailableSlots: (doctorId) => 
    apiRequest(`/appointments/slots/available?doctorUserId=${doctorId}`),
  createAppointment: (data) => 
    apiRequest('/appointments', 'POST', data),
  getPendingAppointments: () => 
    apiRequest('/appointments/pending'),
};

// Payment APIs
export const paymentAPI = {
  createPaymentIntent: (data) => 
    apiRequest('/payments/create-intent', 'POST', data),
  confirmPayment: (data) => 
    apiRequest('/payments/confirm', 'POST', data),
};"""
    
    add_code_box(slide, code, Inches(0.5), Inches(2.2), Inches(9), Inches(4.5), code_bg, text_color)
    
    return slide

def add_security_slide(prs, title_color, text_color, accent_color):
    """Slide 17: Security Implementation"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Security Implementation", title_color)
    
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(5))
    frame = content_box.text_frame
    frame.word_wrap = True
    
    security_features = [
        ("1. Authentication & Authorization", [
            "• JWT token-based authentication with 2-hour expiration",
            "• BCrypt password hashing with configurable work factor",
            "• Role-based access control (PATIENT, DOCTOR)",
            "• Protected routes with [Authorize] attribute"
        ]),
        ("2. Input Validation & Sanitization", [
            "• DTO validation with data annotations",
            "• File type and size validation (10 MB limit)",
            "• SQL injection prevention via EF Core parameterized queries",
            "• XSS prevention through input sanitization"
        ]),
        ("3. API Security", [
            "• HTTPS enforcement with redirect middleware",
            "• Rate limiting on login endpoint (5 attempts/60 seconds)",
            "• CORS configuration (restrict in production)",
            "• Secure error handling (no stack traces exposed)"
        ]),
        ("4. Data Protection", [
            "• Sensitive data in appsettings.json (use Azure Key Vault in production)",
            "• Stripe keys separated (publishable vs. secret)",
            "• Password never stored in plaintext",
            "• JWT signed with secret key"
        ]),
        ("5. File Upload Security", [
            "• Whitelist of allowed MIME types",
            "• Path sanitization to prevent directory traversal",
            "• Unique filename generation (GUID prefix)",
            "• Files stored outside webroot"
        ])
    ]
    
    first = True
    for title, points in security_features:
        p = frame.paragraphs[0] if first else frame.add_paragraph()
        p.text = title
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = accent_color
        if not first:
            p.space_before = Pt(10)
        first = False
        
        for point in points:
            p2 = frame.add_paragraph()
            p2.text = point
            p2.font.size = Pt(10)
            p2.font.color.rgb = text_color
            p2.space_before = Pt(2)
            p2.level = 0
    
    return slide

def add_middleware_slide(prs, title_color, text_color, code_bg):
    """Slide 18: Middleware Pipeline"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Implementation: Middleware Pipeline", title_color)
    
    desc_box = slide.shapes.add_textbox(Inches(1), Inches(1.6), Inches(8), Inches(0.5))
    desc_frame = desc_box.text_frame
    p = desc_frame.paragraphs[0]
    p.text = "Custom middleware for cross-cutting concerns (logging, error handling, authentication)"
    p.font.size = Pt(11)
    p.font.color.rgb = text_color
    
    code = """// Middleware/GlobalExceptionMiddleware.cs - Centralized Error Handling

public class GlobalExceptionMiddleware
{
    private readonly RequestDelegate _next;
    private readonly ILogger<GlobalExceptionMiddleware> _logger;

    public GlobalExceptionMiddleware(
        RequestDelegate next, 
        ILogger<GlobalExceptionMiddleware> logger)
    {
        _next = next;
        _logger = logger;
    }

    public async Task InvokeAsync(HttpContext context)
    {
        try
        {
            await _next(context);
        }
        catch (AppException ex)
        {
            // Known application exception - return friendly error
            _logger.LogWarning(ex, "Application exception occurred");
            
            context.Response.StatusCode = ex.StatusCode;
            context.Response.ContentType = "application/json";
            
            await context.Response.WriteAsJsonAsync(new 
            { 
                error = ex.Message 
            });
        }
        catch (Exception ex)
        {
            // Unknown exception - log details, return generic error
            _logger.LogError(ex, "Unexpected error occurred");
            
            context.Response.StatusCode = 500;
            context.Response.ContentType = "application/json";
            
            await context.Response.WriteAsJsonAsync(new 
            { 
                error = "Internal server error" 
            });
        }
    }
}

// Program.cs - Register middleware
app.UseMiddleware<CorrelationIdMiddleware>();      // 1. Generate request ID
app.UseMiddleware<RequestLoggingMiddleware>();     // 2. Log requests
app.UseMiddleware<GlobalExceptionMiddleware>();    // 3. Handle errors
app.UseAuthentication();                           // 4. Validate JWT
app.UseMiddleware<UserContextMiddleware>();        // 5. Extract user claims"""
    
    add_code_box(slide, code, Inches(0.4), Inches(2.2), Inches(9.2), Inches(4.5), code_bg, text_color)
    
    return slide

def add_repository_slide(prs, title_color, text_color, code_bg):
    """Slide 19: Repository Pattern"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Implementation: Repository Pattern", title_color)
    
    desc_box = slide.shapes.add_textbox(Inches(1), Inches(1.6), Inches(8), Inches(0.5))
    desc_frame = desc_box.text_frame
    p = desc_frame.paragraphs[0]
    p.text = "Data access abstraction layer using Entity Framework Core"
    p.font.size = Pt(11)
    p.font.color.rgb = text_color
    
    code = """// Repositories/UserRepository.cs - Data Access Layer

public class UserRepository : IUserRepository
{
    private readonly ApplicationDbContext _context;

    public UserRepository(ApplicationDbContext context)
    {
        _context = context;
    }

    public async Task<TBL01?> FindByIdAsync(int userId)
    {
        return await _context.Users
            .FirstOrDefaultAsync(u => u.L01F01 == userId);
    }

    public async Task<TBL01?> FindByEmailAsync(string email)
    {
        return await _context.Users
            .FirstOrDefaultAsync(u => u.L01F04 == email);
    }

    public async Task<TBL01> AddAsync(TBL01 user)
    {
        await _context.Users.AddAsync(user);
        await _context.SaveChangesAsync();
        return user;
    }

    public async Task UpdateAsync(TBL01 user)
    {
        _context.Users.Update(user);
        await _context.SaveChangesAsync();
    }

    public async Task<bool> EmailExistsAsync(string email)
    {
        return await _context.Users
            .AnyAsync(u => u.L01F04 == email);
    }
}

// Program.cs - Dependency Injection Registration
services.AddScoped<IUserRepository, UserRepository>();
services.AddScoped<IAppointmentRepository, AppointmentRepository>();
services.AddScoped<IPaymentRepository, PaymentRepository>();"""
    
    add_code_box(slide, code, Inches(0.7), Inches(2.2), Inches(8.6), Inches(4.4), code_bg, text_color)
    
    return slide

def add_api_endpoints_slide(prs, title_color, text_color, accent_color):
    """Slide 20: API Endpoints Overview"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "API Endpoints Overview", title_color)
    
    content_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.7), Inches(8.8), Inches(5.2))
    frame = content_box.text_frame
    frame.word_wrap = True
    
    endpoints = [
        ("Authentication (/api/auth)", [
            "POST /signup - Register new user (Patient/Doctor)",
            "POST /login - User authentication with JWT"
        ]),
        ("Appointments (/api/appointments)", [
            "GET /doctors/available - List doctors with distance",
            "GET /slots/available - Available slots for doctor",
            "POST / - Create appointment request (Patient)",
            "GET /pending, /approved, /completed - Filtered appointments",
            "PUT /{id}/decision - Approve/Decline (Doctor)",
            "PUT /{id}/complete - Mark completed (Doctor)",
            "POST /doctor/slots - Create time slots (Doctor)"
        ]),
        ("Payments (/api/payments)", [
            "POST /create-intent - Create Stripe payment intent",
            "POST /confirm - Confirm payment & create appointment",
            "GET /my-payments - Payment history",
            "POST /refund/{id} - Process refund (Doctor)",
            "GET /earnings - Doctor earnings dashboard"
        ]),
        ("Prescriptions (/api/prescriptions)", [
            "POST / - Create prescription (Doctor)",
            "GET / - Get user prescriptions",
            "GET /{id} - Get specific prescription"
        ]),
        ("Medical Documents (/api/medical-documents)", [
            "POST / - Upload document (Patient)",
            "GET /patient/my-documents - Patient's documents",
            "GET /{id}/download - Download document",
            "DELETE /{id} - Delete document"
        ]),
        ("Doctor & Patient (/api/doctor, /api/patient)", [
            "GET /profile - Get user profile",
            "PUT /profile - Update profile",
            "POST /clinic - Create & associate clinic (Doctor)"
        ])
    ]
    
    first = True
    for category, endpoint_list in endpoints:
        p = frame.paragraphs[0] if first else frame.add_paragraph()
        p.text = category
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = accent_color
        if not first:
            p.space_before = Pt(8)
        first = False
        
        for endpoint in endpoint_list:
            p2 = frame.add_paragraph()
            p2.text = f"• {endpoint}"
            p2.font.size = Pt(9)
            p2.font.color.rgb = text_color
            p2.space_before = Pt(2)
    
    # Total count
    total_box = slide.shapes.add_textbox(Inches(1), Inches(6.8), Inches(8), Inches(0.4))
    total_frame = total_box.text_frame
    p = total_frame.paragraphs[0]
    p.text = "Total Endpoints: 40+ | Documentation: Swagger UI at /swagger"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = title_color
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_features_summary_slide(prs, title_color, text_color, accent_color):
    """Slide 21: Key Features Summary"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Key Features Implemented", title_color)
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(4.2), Inches(5))
    frame_left = left_box.text_frame
    frame_left.word_wrap = True
    
    p1 = frame_left.paragraphs[0]
    p1.text = "Patient Features"
    p1.font.size = Pt(14)
    p1.font.bold = True
    p1.font.color.rgb = accent_color
    
    patient_features = [
        "✓ User registration & authentication",
        "✓ Profile management with health data",
        "✓ Location-based doctor search",
        "✓ Interactive map view (React Leaflet)",
        "✓ Real-time slot availability",
        "✓ Integrated payment (Stripe)",
        "✓ Appointment status tracking",
        "✓ Medical document upload",
        "✓ Digital prescription access",
        "✓ Payment history & billing"
    ]
    
    for feature in patient_features:
        p = frame_left.add_paragraph()
        p.text = feature
        p.font.size = Pt(10)
        p.font.color.rgb = text_color
        p.space_before = Pt(4)
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.8), Inches(4.2), Inches(5))
    frame_right = right_box.text_frame
    frame_right.word_wrap = True
    
    p1 = frame_right.paragraphs[0]
    p1.text = "Doctor Features"
    p1.font.size = Pt(14)
    p1.font.bold = True
    p1.font.color.rgb = accent_color
    
    doctor_features = [
        "✓ Professional profile setup",
        "✓ Multi-clinic management",
        "✓ Appointment slot creation",
        "✓ Request approval/decline",
        "✓ Patient medical records access",
        "✓ Digital prescription creation",
        "✓ Appointment completion tracking",
        "✓ Automatic refund processing",
        "✓ Earnings dashboard",
        "✓ Transaction history"
    ]
    
    for feature in doctor_features:
        p = frame_right.add_paragraph()
        p.text = feature
        p.font.size = Pt(10)
        p.font.color.rgb = text_color
        p.space_before = Pt(4)
    
    return slide

def add_screenshots_slide(prs, title_color, text_color):
    """Slide 22: Screenshots Placeholder"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Application Screenshots", title_color)
    
    # Placeholder boxes for screenshots
    placeholder1 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(2), Inches(4), Inches(2.5)
    )
    placeholder1.fill.solid()
    placeholder1.fill.fore_color.rgb = RGBColor(240, 240, 240)
    placeholder1.line.color.rgb = RGBColor(200, 200, 200)
    
    text1 = placeholder1.text_frame
    text1.text = "[Landing Page Screenshot]"
    text1.paragraphs[0].alignment = PP_ALIGN.CENTER
    text1.paragraphs[0].font.size = Pt(12)
    text1.paragraphs[0].font.color.rgb = RGBColor(150, 150, 150)
    text1.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    placeholder2 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.2), Inches(2), Inches(4), Inches(2.5)
    )
    placeholder2.fill.solid()
    placeholder2.fill.fore_color.rgb = RGBColor(240, 240, 240)
    placeholder2.line.color.rgb = RGBColor(200, 200, 200)
    
    text2 = placeholder2.text_frame
    text2.text = "[Patient Dashboard Screenshot]"
    text2.paragraphs[0].alignment = PP_ALIGN.CENTER
    text2.paragraphs[0].font.size = Pt(12)
    text2.paragraphs[0].font.color.rgb = RGBColor(150, 150, 150)
    text2.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    placeholder3 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(4.7), Inches(4), Inches(2.5)
    )
    placeholder3.fill.solid()
    placeholder3.fill.fore_color.rgb = RGBColor(240, 240, 240)
    placeholder3.line.color.rgb = RGBColor(200, 200, 200)
    
    text3 = placeholder3.text_frame
    text3.text = "[Doctor Dashboard Screenshot]"
    text3.paragraphs[0].alignment = PP_ALIGN.CENTER
    text3.paragraphs[0].font.size = Pt(12)
    text3.paragraphs[0].font.color.rgb = RGBColor(150, 150, 150)
    text3.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    placeholder4 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.2), Inches(4.7), Inches(4), Inches(2.5)
    )
    placeholder4.fill.solid()
    placeholder4.fill.fore_color.rgb = RGBColor(240, 240, 240)
    placeholder4.line.color.rgb = RGBColor(200, 200, 200)
    
    text4 = placeholder4.text_frame
    text4.text = "[Appointment Booking Screenshot]"
    text4.paragraphs[0].alignment = PP_ALIGN.CENTER
    text4.paragraphs[0].font.size = Pt(12)
    text4.paragraphs[0].font.color.rgb = RGBColor(150, 150, 150)
    text4.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Instruction
    instr_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(0.3))
    instr_frame = instr_box.text_frame
    p = instr_frame.paragraphs[0]
    p.text = "Replace placeholder boxes with actual application screenshots"
    p.font.size = Pt(10)
    p.font.italic = True
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_conclusion_slide(prs, title_color, text_color, accent_color):
    """Slide 23: Conclusion"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Conclusion & Future Scope", title_color)
    
    content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(5))
    frame = content_box.text_frame
    frame.word_wrap = True
    
    p1 = frame.paragraphs[0]
    p1.text = "Project Achievements"
    p1.font.size = Pt(14)
    p1.font.bold = True
    p1.font.color.rgb = accent_color
    
    achievements = [
        "✓ Full-stack healthcare appointment system with 40+ API endpoints",
        "✓ Secure authentication using JWT and BCrypt password hashing",
        "✓ Integrated Stripe payment gateway with automatic refund processing",
        "✓ Location-based doctor search using Haversine formula",
        "✓ Digital prescription and medical document management",
        "✓ Clean architecture with Repository and Service patterns",
        "✓ Responsive React frontend with Redux state management",
        "✓ Comprehensive database design with 11 normalized tables"
    ]
    
    for achievement in achievements:
        p = frame.add_paragraph()
        p.text = achievement
        p.font.size = Pt(10)
        p.font.color.rgb = text_color
        p.space_before = Pt(4)
    
    p2 = frame.add_paragraph()
    p2.text = "\nFuture Enhancements"
    p2.font.size = Pt(14)
    p2.font.bold = True
    p2.font.color.rgb = accent_color
    p2.space_before = Pt(12)
    
    future = [
        "• Video consultation integration (Twilio/Agora)",
        "• Email/SMS notification system",
        "• Patient review and rating system",
        "• AI-powered doctor recommendations",
        "• Mobile application (React Native/Flutter)",
        "• Insurance integration",
        "• Advanced analytics dashboard"
    ]
    
    for item in future:
        p = frame.add_paragraph()
        p.text = item
        p.font.size = Pt(10)
        p.font.color.rgb = text_color
        p.space_before = Pt(4)
    
    return slide

def add_references_slide(prs, title_color, text_color):
    """Slide 24: References"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "References", title_color)
    
    content_box = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(7), Inches(4.5))
    frame = content_box.text_frame
    frame.word_wrap = True
    
    references = [
        "[1] Microsoft. (2024). ASP.NET Core Documentation. https://docs.microsoft.com/aspnet/core",
        "[2] React Team. (2024). React Documentation. https://react.dev",
        "[3] Stripe. (2024). Stripe API Documentation. https://stripe.com/docs/api",
        "[4] Entity Framework Core. (2024). EF Core Documentation. https://docs.microsoft.com/ef/core",
        "[5] Redux Toolkit. (2024). Redux Documentation. https://redux-toolkit.js.org",
        "[6] Tailwind Labs. (2024). Tailwind CSS Documentation. https://tailwindcss.com/docs",
        "[7] Leaflet. (2024). React Leaflet Documentation. https://react-leaflet.js.org",
        "[8] JSON Web Tokens. (2024). JWT Introduction. https://jwt.io/introduction",
        "[9] BCrypt. (2024). BCrypt.Net Documentation. https://github.com/BcryptNet/bcrypt.net",
        "[10] MySQL. (2024). MySQL 8.0 Reference Manual. https://dev.mysql.com/doc",
        "",
        "[Add project-specific references as needed]"
    ]
    
    for i, ref in enumerate(references):
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.text = ref
        p.font.size = Pt(9)
        p.font.color.rgb = text_color
        p.space_before = Pt(6)
        if ref.startswith("[Add"):
            p.font.italic = True
            p.font.color.rgb = RGBColor(100, 100, 100)
    
    return slide

def add_slide_title(slide, title_text, color):
    """Add title to slide"""
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title_text
    
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = color
    title_para.alignment = PP_ALIGN.LEFT
    
    # Add horizontal line under title
    line = slide.shapes.add_connector(
        1,  # Straight line
        Inches(0.5), Inches(1.2),
        Inches(9.5), Inches(1.2)
    )
    line.line.color.rgb = RGBColor(68, 114, 196)
    line.line.width = Pt(2)

def add_code_box(slide, code, left, top, width, height, bg_color, text_color):
    """Add code snippet box"""
    # Background box
    code_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    code_bg.fill.solid()
    code_bg.fill.fore_color.rgb = bg_color
    code_bg.line.color.rgb = RGBColor(220, 220, 220)
    
    # Code text
    code_box = slide.shapes.add_textbox(
        left + Inches(0.15), 
        top + Inches(0.1), 
        width - Inches(0.3), 
        height - Inches(0.2)
    )
    frame = code_box.text_frame
    frame.word_wrap = False
    frame.text = code
    
    para = frame.paragraphs[0]
    para.font.name = 'Courier New'
    para.font.size = Pt(8)
    para.font.color.rgb = text_color

if __name__ == "__main__":
    create_presentation()
